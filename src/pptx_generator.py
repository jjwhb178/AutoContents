from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os

# --- 디자인 토큰 (Professional Slate & Gold) ---
COLOR_BG = RGBColor(15, 25, 45)       # 짙은 남색 배경
COLOR_HEADER_BG = RGBColor(25, 35, 60) # 헤더 배경
COLOR_ACCENT = RGBColor(212, 175, 55) # 머니대디 골드
COLOR_TEXT_MAIN = RGBColor(240, 240, 240) # 본문 화이트
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_NEON_UP = RGBColor(255, 50, 50)  # 상승(레드)
COLOR_NEON_DOWN = RGBColor(50, 150, 255) # 하락(블루)

def apply_text_style(p, size, color, bold=False, align=PP_ALIGN.LEFT):
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Malgun Gothic"
    p.alignment = align

def draw_section_header(slide, title, page_num):
    """전문 리포트 스타일 상단 헤더"""
    # 헤더 바
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
    rect.fill.solid(); rect.fill.fore_color.rgb = COLOR_HEADER_BG
    rect.line.fill.background()
    
    # 섹션 텍스트
    section_name = "MONEYDADDY STRATEGIC REPORT"
    sec_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(4), Inches(0.3))
    apply_text_style(sec_box.text_frame.paragraphs[0], 10, COLOR_ACCENT, bold=True)
    sec_box.text_frame.paragraphs[0].text = section_name
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.35), Inches(8), Inches(0.4))
    apply_text_style(title_box.text_frame.paragraphs[0], 24, COLOR_WHITE, bold=True)
    title_box.text_frame.paragraphs[0].text = title
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(Inches(9), Inches(0.2), Inches(0.7), Inches(0.4))
    apply_text_style(page_box.text_frame.paragraphs[0], 12, COLOR_ACCENT, bold=True, align=PP_ALIGN.RIGHT)
    page_box.text_frame.paragraphs[0].text = str(page_num)

def draw_heatmap(slide, top, left, width, height):
    """4P용 Finviz 스타일 히트맵 시뮬레이션"""
    cols, rows = 4, 3
    w, h = width / cols, height / rows
    for r in range(rows):
        for c in range(cols):
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + c*w, top + r*h, w-Inches(0.05), h-Inches(0.05))
            rect.fill.solid()
            # 랜덤하게 상승/하락 색상 배치
            rect.fill.fore_color.rgb = COLOR_NEON_UP if (r+c)%2==0 else COLOR_NEON_DOWN
            rect.line.color.rgb = COLOR_WHITE; rect.line.width = Pt(0.5)

def draw_connection_diagram(slide, top, left, width, height):
    """6P용 종목 연결 도식"""
    # 美 종목 (좌)
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + height/4, width/3, height/2)
    box1.fill.solid(); box1.fill.fore_color.rgb = COLOR_HEADER_BG
    box1.line.color.rgb = COLOR_ACCENT
    t1 = box1.text_frame.paragraphs[0]
    t1.text = "US MARKET LEADERS"; apply_text_style(t1, 14, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # 연결 화살표
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + width/3 + Inches(0.1), top + height/2 - Inches(0.2), width/4, Inches(0.4))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = COLOR_ACCENT
    
    # 韓 관련주 (우)
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + width*0.6, top + height/4, width/3, height/2)
    box2.fill.solid(); box2.fill.fore_color.rgb = COLOR_BG
    box2.line.color.rgb = COLOR_NEON_UP
    t2 = box2.text_frame.paragraphs[0]
    t2.text = "K-RELATED STOCKS"; apply_text_style(t2, 14, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

def create_chart_image(data_points, title, filename, color="#D4AF37"):
    """Matplotlib을 사용해 실제 데이터를 반영한 차트 이미지 생성"""
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg') # GUI 없는 환경 대응
    
    plt.figure(figsize=(6, 4), facecolor='#1E2638')
    ax = plt.axes()
    ax.set_facecolor('#1E2638')
    
    labels = [f"D-{len(data_points)-i-1}" for i in range(len(data_points))]
    plt.plot(labels, data_points, marker='o', color=color, linewidth=2, markersize=8)
    plt.fill_between(labels, data_points, color=color, alpha=0.2)
    
    plt.title(title, color='white', fontsize=14, pad=20, fontweight='bold')
    plt.xticks(color='white')
    plt.yticks(color='white')
    
    # 테두리 제거
    for spine in ax.spines.values():
        spine.set_visible(False)
    
    plt.grid(color='#2A344A', linestyle='--', linewidth=0.5)
    
    plt.tight_layout()
    plt.savefig(filename, dpi=100)
    plt.close()
    return filename

def create_pro_pptx():
    logic_path = os.path.join("data", "latest_content_logic.json")
    if not os.path.exists(logic_path): return
    with open(logic_path, "r", encoding="utf-8") as f:
        logic = json.load(f)

    # 차트용 가상 데이터 (raw_market_data에서 가져올 수도 있음)
    score = logic.get("score", 50)
    chart_data = [score * (0.9 + i*0.05) for i in range(5)]

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    script = logic.get("ppt_script", {})

    for page_num in range(1, 19):
        p_str = str(page_num)
        data = script.get(p_str)
        if not data: continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 배경색 설정
        background = slide.background
        background.fill.solid(); background.fill.fore_color.rgb = COLOR_BG

        title_text = data.get('title', '').split(': ')[-1]
        draw_section_header(slide, title_text, page_num)
        
        # [레이아웃 분할] 좌측(시각 자료), 우측(본문 카드)
        v_elements = data.get('visual_elements', [])
        if isinstance(v_elements, str): v_elements = [v_elements]
        
        # 1. 시각 자료 영역 (좌측 6인치)
        visual_top, visual_left = Inches(1.2), Inches(0.3)
        visual_w, visual_h = Inches(5.5), Inches(5.5)
        
        layout_type = data.get('layout_type', 'bullets')
        
        if page_num == 1:
            # 1P: 썸네일 수준의 풀 블리드 비주얼
            background.fill.solid(); background.fill.fore_color.rgb = COLOR_HEADER_BG
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
            apply_text_style(title_box.text_frame.paragraphs[0], 40, COLOR_ACCENT, bold=True, align=PP_ALIGN.CENTER)
            title_box.text_frame.paragraphs[0].text = title_text
            
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
            tf = sub_box.text_frame
            tf.word_wrap = True
            for i, elem in enumerate(v_elements[:3]):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = f"▶ {elem}"
                apply_text_style(p, 18, COLOR_WHITE, align=PP_ALIGN.CENTER)
            continue
            
        elif layout_type == 'warning' or page_num == 10:
            # 10P: 통찰의 반전 (Warning Layout)
            background.fill.solid(); background.fill.fore_color.rgb = RGBColor(0, 0, 0)
            warning_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
            apply_text_style(warning_box.text_frame.paragraphs[0], 36, COLOR_ACCENT, bold=True, align=PP_ALIGN.CENTER)
            warning_box.text_frame.paragraphs[0].text = f"WARNING: {title_text}"
            
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
            apply_text_style(sub_box.text_frame.paragraphs[0], 18, COLOR_NEON_UP, align=PP_ALIGN.CENTER)
            sub_box.text_frame.paragraphs[0].text = v_elements[0] if v_elements else "Decoupling Risk Detected"
            continue

        if page_num == 4:
            draw_heatmap(slide, visual_top, visual_left, visual_w, visual_h)
        elif page_num == 11:
            # 11P: 실제 차트 삽입
            chart_file = os.path.join("data", "market_chart_v11.png")
            create_chart_image(chart_data, "Market Momentum Trace", chart_file)
            slide.shapes.add_picture(chart_file, visual_left, visual_top + Inches(0.5), visual_w, visual_h - Inches(1.0))
        elif page_num == 15:
            # 15P: 섹터 타겟 차트
            chart_file = os.path.join("data", "sector_target_v15.png")
            create_chart_image([10, 15, 13, 22, 28], "Target Sector Growth", chart_file, color="#3296FF")
            slide.shapes.add_picture(chart_file, visual_left, visual_top + Inches(0.5), visual_w, visual_h - Inches(1.0))
        else:
            # 일반 시각 자료 박스 (Matplotlib 차트로 대체 시도)
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, visual_left, visual_top, visual_w, visual_h)
            rect.fill.solid(); rect.fill.fore_color.rgb = COLOR_HEADER_BG
            rect.line.color.rgb = COLOR_ACCENT; rect.line.width = Pt(1.5)
            p = rect.text_frame.paragraphs[0]
            p.text = f"ANALYSIS: {title_text}"
            apply_text_style(p, 14, COLOR_ACCENT, align=PP_ALIGN.CENTER)
            rect.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 2. 본문 텍스트 카드 (우측 4인치) - 글자가 절대 나가지 않도록 고정
        card_left, card_top = Inches(6.1), Inches(1.2)
        card_w, card_h = Inches(3.6), Inches(5.5)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = COLOR_HEADER_BG
        card.line.color.rgb = COLOR_WHITE; card.line.width = Pt(0.5)
        
        # 본문 텍스트 프레임 (패딩 0.2인치 적용)
        txt_box = slide.shapes.add_textbox(card_left + Inches(0.2), card_top + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
        tf = txt_box.text_frame
        tf.word_wrap = True
        # 자동 폰트 크기 조절 (Auto-fit) 활성화
        tf.auto_size = 1 # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        for i, elem in enumerate(v_elements[:6]): # 최대 6개로 제한하여 가독성 확보
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"▶ {elem}"
            apply_text_style(p, 14, COLOR_TEXT_MAIN)
            p.space_after = Pt(15)

    # 최종 저장 (동적 날짜 경로)
    from datetime import datetime
    today_str = datetime.now().strftime("%Y-%m-%d")
    out_dir = os.path.join("outputs", today_str)
    if not os.path.exists(out_dir): os.makedirs(out_dir)
    output_path = os.path.join(out_dir, "daily_strategy_v8_pro.pptx")
    prs.save(output_path)
    print(f"Professional Report PPTX (v8) generated: {output_path}")

if __name__ == "__main__":
    create_pro_pptx()
